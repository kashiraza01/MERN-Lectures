"""
Generate Day 26 — Node.js Basics Deep Dive & Recap (90-minute lecture deck).
Run: python generate_node_basics_lecture.py
Output: Lecture Day 26 - Node.js Basics Recap.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT = "Lecture Day 26 - Node.js Basics Recap.pptx"

GREEN = RGBColor(0x68, 0xA0, 0x63)  # Node green
DARK = RGBColor(0x23, 0x2B, 0x2B)
GRAY = RGBColor(0x5C, 0x6C, 0x6C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0xF4, 0xF4, 0xF4)


def set_notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = GRAY
        p2.space_before = Pt(12)
    return slide


def add_section_slide(prs, section_title, duration=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    if duration:
        p2 = tf.add_paragraph()
        p2.text = duration
        p2.font.size = Pt(20)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)
    return slide


def add_bullet_slide(prs, title, bullets, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.8))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = DARK
        p.level = 0
        p.space_after = Pt(10)

    if notes:
        set_notes(slide, notes)
    return slide


def add_code_slide(prs, title, code, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    code_box = slide.shapes.add_shape(1, Inches(0.6), Inches(1.2), Inches(8.8), Inches(5.8))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = CODE_BG
    code_box.line.color.rgb = GRAY

    tf = code_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = code.strip()
    p.font.name = "Consolas"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK

    if notes:
        set_notes(slide, notes)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    s = add_title_slide(
        prs,
        "Node.js Basics — Deep Dive & Recap",
        "Week 6 · Day 26 · 90 minutes",
    )
    set_notes(
        s,
        "Welcome. Students already saw Express + Mongoose (Day 23) and a Node intro (Day 24). "
        "Today we recap Node fundamentals thoroughly so the backend stack makes sense.",
    )

    add_bullet_slide(
        prs,
        "Agenda (90 min)",
        [
            "Where Node fits in MERN (5 min)",
            "What Node is & how it runs JavaScript (10 min)",
            "Browser vs Node runtime (8 min)",
            "Running scripts, REPL & npm (10 min)",
            "Modules: CommonJS & ES Modules (12 min)",
            "Core modules: path, fs, os, process (15 min)",
            "Event loop & async I/O (15 min)",
            "Events, streams & raw HTTP (12 min)",
            "package.json, env vars & project layout (8 min)",
            "Bridge to Express/Mongoose & lab preview (5 min)",
        ],
    )

    # Section 1
    add_section_slide(prs, "Section 1 — Node in the MERN Stack", "5 min")

    add_bullet_slide(
        prs,
        "Recap: what you already built",
        [
            "Day 22 — MongoDB: document database, collections, CRUD in shell/Compass",
            "Day 23 — Express + Mongoose: REST API, routes, schemas, models",
            "Day 24 — Node introduction: JavaScript on the server",
            "Today — understand the engine under Express: Node.js itself",
            "MERN = MongoDB + Express + React + Node",
        ],
        notes="Draw the stack: React (browser) → HTTP → Express (Node) → Mongoose → MongoDB.",
    )

    add_bullet_slide(
        prs,
        "Why Node for the backend?",
        [
            "Same language (JavaScript) on frontend and backend",
            "Huge npm ecosystem — Express, Mongoose, dotenv, etc.",
            "Non-blocking I/O — handles many concurrent connections efficiently",
            "Event-driven architecture fits real-time & API workloads",
            "Express is a thin layer on top of Node's built-in http module",
        ],
    )

    # Section 2
    add_section_slide(prs, "Section 2 — What Is Node.js?", "10 min")

    add_bullet_slide(
        prs,
        "Node.js in one sentence",
        [
            "Node.js = JavaScript runtime built on Chrome's V8 engine + libuv",
            "V8 compiles JS to machine code (fast execution)",
            "libuv provides the event loop, thread pool & async I/O",
            "Node is NOT a framework — it's a runtime environment",
            "Express, Fastify, NestJS are frameworks that run ON Node",
        ],
        notes="Analogy: Node is the engine; Express is the car body.",
    )

    add_bullet_slide(
        prs,
        "Key characteristics",
        [
            "Single-threaded for JavaScript execution (one call stack)",
            "Non-blocking I/O — async operations don't freeze the server",
            "Cross-platform — Windows, macOS, Linux",
            "Open source, maintained by the OpenJS Foundation",
            "LTS releases — use even-numbered versions in production (e.g. 22.x)",
        ],
    )

    add_code_slide(
        prs,
        "Your first Node program",
        """// hello.js
console.log("Hello from Node!");
console.log("Node version:", process.version);
console.log("Platform:", process.platform);""",
        notes="Run: node hello.js. No browser, no HTML — just JS + Node APIs.",
    )

    # Section 3
    add_section_slide(prs, "Section 3 — Browser vs Node", "8 min")

    add_bullet_slide(
        prs,
        "Globals: browser vs Node",
        [
            "Browser: window, document, localStorage, fetch (built-in)",
            "Node: global (or globalThis), process, __dirname, __filename",
            "Both: console, setTimeout, Promise, JSON, Map, Set",
            "Node has NO DOM — no document.getElementById",
            "Node adds file system, networking, child processes via modules",
        ],
    )

    add_bullet_slide(
        prs,
        "APIs you will use constantly",
        [
            "fs — read/write files",
            "path — join paths safely across OS",
            "http / https — create servers & make requests",
            "os — CPU, memory, hostname info",
            "process — env vars, argv, exit codes",
            "events — EventEmitter pattern (used inside streams & servers)",
        ],
    )

    # Section 4
    add_section_slide(prs, "Section 4 — Running Node & npm", "10 min")

    add_bullet_slide(
        prs,
        "Three ways to run JavaScript with Node",
        [
            "1. node script.js — run a file",
            "2. node — interactive REPL (Read-Eval-Print Loop)",
            "3. npm run <script> — run commands from package.json",
            "node --watch app.js — auto-restart on file changes (Node 18+)",
            "npx nodemon app.js — popular dev tool for auto-restart",
        ],
    )

    add_code_slide(
        prs,
        "REPL quick demo",
        """$ node
> 2 + 2
4
> const fs = require('fs')
> fs.readdirSync('.')
[ 'package.json', 'src', 'node_modules' ]
> .exit""",
        notes="REPL is great for quick experiments. .help shows commands.",
    )

    add_bullet_slide(
        prs,
        "npm essentials (recap)",
        [
            "npm init -y — create package.json",
            "npm install express — add dependency (saves to node_modules/)",
            "npm install -D nodemon — dev dependency",
            "package-lock.json — locks exact versions for reproducible installs",
            "Never commit node_modules/ — always use .gitignore",
        ],
    )

    # Section 5
    add_section_slide(prs, "Section 5 — Modules", "12 min")

    add_bullet_slide(
        prs,
        "Why modules?",
        [
            "Split code into reusable files",
            "Avoid global namespace pollution",
            "You learned ES modules in Week 2 (import/export with Vite)",
            "Node supports BOTH CommonJS and ES Modules",
            "Express projects often mix: require() in older code, import in newer",
        ],
    )

    add_code_slide(
        prs,
        "CommonJS (require / module.exports)",
        """// math.js
function add(a, b) { return a + b; }
function subtract(a, b) { return a - b; }
module.exports = { add, subtract };

// app.js
const { add, subtract } = require('./math');
console.log(add(5, 3));       // 8
console.log(subtract(5, 3));  // 2""",
    )

    add_code_slide(
        prs,
        "ES Modules (import / export)",
        """// math.mjs  OR  set "type": "module" in package.json
export function add(a, b) { return a + b; }
export function subtract(a, b) { return a - b; }

// app.mjs
import { add, subtract } from './math.mjs';
console.log(add(5, 3));""",
        notes='With "type": "module", .js files use import/export. __dirname is not available — use import.meta.url.',
    )

    add_bullet_slide(
        prs,
        "Default vs named exports",
        [
            "Named: export const foo = ... → import { foo } from './file'",
            "Default: export default function() {} → import myFn from './file'",
            "You can mix one default + many named exports in one file",
            "require() is synchronous; import() can be dynamic: await import('./x')",
            "Rule of thumb: pick one style per project and stay consistent",
        ],
    )

    # Section 6
    add_section_slide(prs, "Section 6 — Core Modules", "15 min")

    add_code_slide(
        prs,
        "path — cross-platform file paths",
        """const path = require('path');

path.join('users', 'ali', 'notes.txt');
// users\\ali\\notes.txt  (Windows)
// users/ali/notes.txt    (macOS/Linux)

path.basename('/data/report.pdf');  // report.pdf
path.extname('photo.jpg');          // .jpg
path.resolve('src', 'app.js');      // absolute path from cwd""",
    )

    add_code_slide(
        prs,
        "fs — read & write files (async vs sync)",
        """const fs = require('fs');

// Async (preferred — non-blocking)
fs.readFile('data.txt', 'utf8', (err, data) => {
  if (err) throw err;
  console.log(data);
});

// Sync (blocks the event loop — avoid in servers)
const data = fs.readFileSync('data.txt', 'utf8');

// Promises API (Node 10+)
const fsp = require('fs/promises');
const text = await fsp.readFile('data.txt', 'utf8');""",
        notes="In Express route handlers, always prefer async fs or streams for large files.",
    )

    add_code_slide(
        prs,
        "process & os",
        """// process — current Node process
process.env.PORT          // environment variables
process.argv              // CLI arguments
process.cwd()             // current working directory
process.exit(1)           // exit with error code

// os — machine info
const os = require('os');
os.platform();            // win32, darwin, linux
os.cpus().length;         // CPU cores
os.totalmem();            // RAM in bytes""",
    )

    add_bullet_slide(
        prs,
        "Environment variables (.env)",
        [
            "Never hard-code secrets (DB passwords, API keys) in source code",
            "Use process.env.MONGO_URI, process.env.PORT, etc.",
            "dotenv package loads .env file into process.env at startup",
            "Add .env to .gitignore — commit .env.example with placeholder values",
            "This is exactly how your Day 23 Express + Mongoose app connects to Atlas",
        ],
    )

    # Section 7
    add_section_slide(prs, "Section 7 — Event Loop & Async I/O", "15 min")

    add_bullet_slide(
        prs,
        "How Node stays fast with one thread",
        [
            "Call stack — runs synchronous JS line by line",
            "Node delegates slow work (disk, network) to libuv thread pool / OS",
            "When async work finishes, callbacks enter the callback queue",
            "Event loop picks queued callbacks when the stack is empty",
            "This is why blocking sync code (fs.readFileSync in a loop) kills performance",
        ],
    )

    add_bullet_slide(
        prs,
        "Microtasks vs macrotasks",
        [
            "Microtasks (highest priority): Promise.then, queueMicrotask",
            "Macrotasks: setTimeout, setInterval, I/O callbacks",
            "Order: sync code → all microtasks → one macrotask → repeat",
            "You practiced this in Week 2 Day 3 lab!",
            "Understanding this prevents subtle bugs in async code",
        ],
    )

    add_code_slide(
        prs,
        "Event loop — predict the output",
        """console.log('A');

setTimeout(() => console.log('B'), 0);

Promise.resolve()
  .then(() => console.log('C'))
  .then(() => console.log('D'));

queueMicrotask(() => console.log('E'));

console.log('F');

// Answer: A, F, C, E, D, B
// Sync first, then microtasks, then macrotasks""",
        notes="Live-poll the class before running. This is Lab Task 1.",
    )

    add_code_slide(
        prs,
        "Async patterns recap",
        """// Callback
fs.readFile('f.txt', 'utf8', (err, data) => { /* ... */ });

// Promise
fsp.readFile('f.txt', 'utf8').then(console.log).catch(console.error);

// async/await (cleanest)
async function load() {
  try {
    const data = await fsp.readFile('f.txt', 'utf8');
    console.log(data);
  } catch (err) {
    console.error(err);
  }
}""",
    )

    # Section 8
    add_section_slide(prs, "Section 8 — Events, Streams & HTTP", "12 min")

    add_code_slide(
        prs,
        "EventEmitter pattern",
        """const EventEmitter = require('events');
const emitter = new EventEmitter();

emitter.on('order', (item) => {
  console.log('Order received:', item);
});

emitter.emit('order', 'Pizza');
// Order received: Pizza

// Many Node objects ARE EventEmitters: streams, servers, process""",
    )

    add_bullet_slide(
        prs,
        "Streams (conceptual)",
        [
            "Streams process data chunk-by-chunk instead of loading entire file into RAM",
            "Types: Readable, Writable, Duplex, Transform",
            "Use case: large file upload/download, log tailing, video piping",
            "fs.createReadStream() + pipe() — classic Node pattern",
            "Express req/res objects are streams under the hood",
        ],
    )

    add_code_slide(
        prs,
        "Raw HTTP server (before Express)",
        """const http = require('http');

const server = http.createServer((req, res) => {
  console.log(req.method, req.url);
  res.writeHead(200, { 'Content-Type': 'application/json' });
  res.end(JSON.stringify({ message: 'Hello from raw Node HTTP!' }));
});

server.listen(3000, () => {
  console.log('Server running on http://localhost:3000');
});""",
        notes="Express app.listen() wraps this. req/res are enhanced versions of these objects.",
    )

    add_bullet_slide(
        prs,
        "Express builds on Node http",
        [
            "const app = express() creates an application object",
            "app.get('/users', handler) registers route + method matching",
            "app.use(middleware) — chain of functions (req, res, next)",
            "Mongoose connect() runs when Node process starts, before routes handle traffic",
            "Full stack: Client → Express routes → Mongoose model → MongoDB",
        ],
    )

    # Section 9
    add_section_slide(prs, "Section 9 — Project Structure & package.json", "8 min")

    add_code_slide(
        prs,
        "Typical Node API project layout",
        """shop-api/
├── package.json
├── .env                 # secrets (gitignored)
├── .env.example         # template for teammates
├── .gitignore
├── server.js            # entry point
├── src/
│   ├── routes/
│   ├── models/          # Mongoose schemas
│   ├── controllers/
│   └── middleware/
└── node_modules/        # gitignored""",
    )

    add_code_slide(
        prs,
        "package.json scripts",
        """{
  "name": "shop-api",
  "version": "1.0.0",
  "main": "server.js",
  "type": "commonjs",
  "scripts": {
    "start": "node server.js",
    "dev": "nodemon server.js",
    "test": "node --test"
  },
  "dependencies": {
    "express": "^4.21.0",
    "mongoose": "^8.0.0",
    "dotenv": "^16.4.0"
  },
  "devDependencies": {
    "nodemon": "^3.1.0"
  }
}""",
    )

    add_bullet_slide(
        prs,
        "Common mistakes to avoid",
        [
            "Using var instead of const/let",
            "Forgetting await on async Mongoose queries",
            "Not handling errors in async route handlers (try/catch or wrapper)",
            "Committing .env or node_modules to GitHub",
            "Blocking the event loop with heavy sync computation in route handlers",
            "Missing process.on('unhandledRejection') logging in production apps",
        ],
    )

    # Wrap up
    add_section_slide(prs, "Section 10 — Summary & Lab", "5 min")

    add_bullet_slide(
        prs,
        "Key takeaways",
        [
            "Node = V8 + libuv; runs JS outside the browser",
            "Core modules (fs, path, http, process) power every Express app",
            "Event loop: sync → microtasks → macrotasks",
            "Modules organize code; CommonJS & ESM both supported",
            "Express/Mongoose from Day 23 sit on top of these fundamentals",
            "Lab today: event loop, fs/path, EventEmitter, raw HTTP server",
        ],
    )

    s = add_title_slide(prs, "Thank you!", "Lab Day 26 · node-basics-lab")
    set_notes(s, "Open lab-starter folder. Students work in pairs for first 15 min, then solo.")

    prs.save(OUTPUT)
    print(f"Saved {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
